import streamlit as st
import boto3
from anthropic import Anthropic
from botocore.config import Config
import shutil
import os
import pandas as pd
import time
import json
import base64
import io
import re
import numpy as np
import openpyxl
from python_calamine import CalamineWorkbook
from openpyxl.cell import Cell
from openpyxl.worksheet.cell_range import CellRange
from boto3.dynamodb.conditions import Key 
import uuid
from pptx import Presentation
from botocore.exceptions import ClientError
from textractor import Textractor
from textractor.visualizers.entitylist import EntityList
from textractor.data.constants import TextractFeatures
from textractor.data.text_linearization_config import TextLinearizationConfig
import pytesseract
from PIL import Image
import PyPDF2
import chardet
import random
from datetime import datetime    
from docx import Document as DocxDocument
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.document import Document
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from docx.table import Table as DocxTable
import concurrent.futures
from functools import partial
import csv
import textract
config = Config(
    read_timeout=600, # Read timeout parameter
    retries = dict(
        max_attempts = 10 ## Handle retries
    )
)

st.set_page_config(initial_sidebar_state="auto")
# Read credentials
with open('config.json','r',encoding='utf-8') as f:
    config_file = json.load(f)
# pricing info
with open('pricing.json','r',encoding='utf-8') as f:
    pricing_file = json.load(f)

S3 = boto3.client('s3')
DYNAMODB  = boto3.resource('dynamodb')

LOCAL_CHAT_FILE_NAME = "chatbot-history.json"
DYNAMODB_TABLE=config_file["DynamodbTable"]
BUCKET=config_file["Bucket_Name"]
OUTPUT_TOKEN=config_file["max-output-token"]
S3_DOC_CACHE_PATH=config_file["document-upload-cache-s3-path"]
TEXTRACT_RESULT_CACHE_PATH=config_file["AmazonTextract-result-cache"]
LOAD_DOC_IN_ALL_CHAT_CONVO=config_file["load-doc-in-chat-history"]
CHAT_HISTORY_LENGTH=config_file["chat-history-loaded-length"]
DYNAMODB_USER=config_file["UserId"]
REGION=config_file["bedrock-region"]
USE_TEXTRACT=config_file["AmazonTextract"]
CSV_SEPERATOR=config_file["csv-delimiter"]

bedrock_runtime = boto3.client(service_name='bedrock-runtime',region_name=REGION,config=config)

if 'messages' not in st.session_state:
    st.session_state['messages'] = []
if 'input_token' not in st.session_state:
    st.session_state['input_token'] = 0
if 'output_token' not in st.session_state:
    st.session_state['output_token'] = 0
if 'chat_hist' not in st.session_state:
    st.session_state['chat_hist'] = []
if 'user_sess' not in st.session_state:
    st.session_state['user_sess'] =str(time.time())
if 'chat_session_list' not in st.session_state:
    st.session_state['chat_session_list'] = []
if 'count' not in st.session_state:
    st.session_state['count'] = 0
if 'userid' not in st.session_state:
    st.session_state['userid']= config_file["UserId"]
if 'cost' not in st.session_state:
    st.session_state['cost'] = 0

def get_object_with_retry(bucket, key):
    max_retries=5
    retries = 0   
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    s3 = boto3.client('s3')
    while retries < max_retries:
        try:
            response = s3.get_object(Bucket=bucket, Key=key)
            return response
        except ClientError as e:
            error_code = e.response['Error']['Code']
            if error_code == 'DecryptionFailureException':
                sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                print(f"Decryption failed, retrying in {sleep_time} seconds...")                
                time.sleep(sleep_time)               
                retries += 1
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
            else:
                raise e

    # If we reach this point, it means the maximum number of retries has been exceeded
    raise Exception(f"Failed to get object {key} from bucket {bucket} after {max_retries} retries.")


def save_chat_local(file_path, new_data, session_id):
    """Store long term chat history Local Disk"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
        if session_id not in existing_data:
            existing_data[session_id]=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = {session_id:[]}
    # Append the new data to the existing list
    from decimal import Decimal
    data = [{k: float(v) if isinstance(v, Decimal) else v for k, v in item.items()} for item in new_data]
    existing_data[session_id].extend(data)

    # Write the updated list back to the JSON file
    with open(file_path, "w",encoding="utf-8") as file:
        json.dump(existing_data, file)
        
def load_chat_local(file_path,session_id):
    """Load long term chat history from Local"""   
    try:
        # Read the existing JSON data from the file
        with open(file_path, "r",encoding='utf-8') as file:
            existing_data = json.load(file)
            if session_id in existing_data:
                existing_data=existing_data[session_id]
            else:
                existing_data=[]
    except FileNotFoundError:
        # If the file doesn't exist, initialize an empty list
        existing_data = []
    return existing_data
    
    
def process_files(files):

    result_string=""
    errors = []
    future_proxy_mapping = {} 
    futures = []

    with concurrent.futures.ProcessPoolExecutor() as executor:
        # Partial function to pass the handle_doc_upload_or_s3 function
        func = partial(handle_doc_upload_or_s3)   
        for file in files:
            future = executor.submit(func, file)
            future_proxy_mapping[future] = file
            futures.append(future)

        # Collect the results and handle exceptions
        for future in concurrent.futures.as_completed(futures):        
            file_url= future_proxy_mapping[future]
            try:
                result = future.result()               
                doc_name=os.path.basename(file_url)
                
                result_string+=f"<{doc_name}>\n{result}\n</{doc_name}>\n"
            except Exception as e:
                # Get the original function arguments from the Future object
                error = {'file': file_url, 'error': str(e)}
                errors.append(error)

    return errors, result_string

def handle_doc_upload_or_s3(file):
    """Handle various document format"""
    dir_name, ext = os.path.splitext(file)
    if  ext.lower() in [".pdf", ".png", ".jpg",".tif",".jpeg"]:   
        content=exract_pdf_text_aws(file)
    elif ".csv"  == ext.lower():
        content=parse_csv_from_s3(file)
    elif ext.lower() in [".xlsx", ".xls"]:
        content=table_parser_utills(file)   
    elif  ".json"==ext.lower():      
        obj=get_s3_obj_from_bucket_(file)
        content = json.loads(obj['Body'].read())  
    elif  ext.lower() in [".txt",".py"]:       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
    elif ".docx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)
        content = extract_text_and_tables(docx_buffer)
    elif ".pptx" == ext.lower():       
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        docx_buffer = io.BytesIO(content)        
        content = extract_text_from_pptx_s3(docx_buffer)
    else:            
        obj=get_s3_obj_from_bucket_(file)
        content = obj['Body'].read()
        doc_buffer = io.BytesIO(content)
        content = textract.process(doc_buffer).decode()
    # Implement any other file extension logic 
    return content

class InvalidContentError(Exception):
    pass

def detect_encoding(s3_uri):
    """detect csv encoding"""
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", s3_uri)
    if match:
        bucket_name = match.group(1)
        key = match.group(2) 
    response = s3.get_object(Bucket=bucket_name, Key=key)
    content = response['Body'].read()
    result = chardet.detect(content)
    return result['encoding']

def parse_csv_from_s3(s3_uri):
    """read csv files"""
    try:
        # Detect the file encoding using chardet
        encoding = detect_encoding(s3_uri)        
        # Sniff the delimiter and read the CSV file
        df = pd.read_csv(s3_uri, delimiter=None, engine='python', encoding=encoding)
        return df.to_csv(index=False, sep=CSV_SEPERATOR)
    except Exception as e:
        raise InvalidContentError(f"Error: {e}")
    
def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

def extract_text_and_tables(docx_path):
    """ Extract text from docx files"""
    document = DocxDocument(docx_path)
    content = ""
    current_section = ""
    section_type = None
    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            if block.text:
                if block.style.name == 'Heading 1':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None  
                    section_type ="h1"
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                elif block.style.name== 'Heading 3':
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                    section_type = "h3"  
                    content += f"<{section_type}>{block.text}</{section_type}>\n"
                
                elif block.style.name == 'List Paragraph':
                    # Add to the current list section
                    if section_type != "list":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "list"
                        current_section = "<list>"
                    current_section += f"{block.text}\n"
                elif block.style.name.startswith('toc'):
                    # Add to the current toc section
                    if section_type != "toc":
                        # Close the current section if it exists
                        if current_section:
                            content += f"{current_section}</{section_type}>\n"
                        section_type = "toc"
                        current_section = "<toc>"
                    current_section += f"{block.text}\n"
                else:
                    # Close the current section if it exists
                    if current_section:
                        content += f"{current_section}</{section_type}>\n"
                        current_section = ""
                        section_type = None
                    
                    # Append the passage text without tagging
                    content += f"{block.text}\n"
        
        elif isinstance(block, DocxTable):
            # Add the current section before the table
            if current_section:
                content += f"{current_section}</{section_type}>\n"
                current_section = ""
                section_type = None

            content += "<table>\n"
            for row in block.rows:
                row_content = []
                for cell in row.cells:
                    cell_content = []
                    for nested_block in iter_block_items(cell):
                        if isinstance(nested_block, Paragraph):
                            cell_content.append(nested_block.text)
                        elif isinstance(nested_block, DocxTable):
                            nested_table_content = parse_nested_table(nested_block)
                            cell_content.append(nested_table_content)
                    row_content.append(CSV_SEPERATOR.join(cell_content))
                content += CSV_SEPERATOR.join(row_content) + "\n"
            content += "</table>\n"

    # Add the final section
    if current_section:
        content += f"{current_section}</{section_type}>\n"

    return content

def parse_nested_table(table):
    nested_table_content = "<table>\n"
    for row in table.rows:
        row_content = []
        for cell in row.cells:
            cell_content = []
            for nested_block in iter_block_items(cell):
                if isinstance(nested_block, Paragraph):
                    cell_content.append(nested_block.text)
                elif isinstance(nested_block, DocxTable):
                    nested_table_content += parse_nested_table(nested_block)
            row_content.append(CSV_SEPERATOR.join(cell_content))
        nested_table_content += CSV_SEPERATOR.join(row_content) + "\n"
    nested_table_content += "</table>"
    return nested_table_content



def extract_text_from_pptx_s3(pptx_buffer):
    """ Extract Text from pptx files"""
    presentation = Presentation(pptx_buffer)    
    text_content = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text_content.append('\n'.join(slide_text))    
    return '\n\n'.join(text_content)
    
def exract_pdf_text_aws(file):    
    file_base_name=os.path.basename(file)
    dir_name, ext = os.path.splitext(file)
    # Checking if extracted doc content is in S3
    if USE_TEXTRACT:        
        if [x for x in get_s3_keys(f"{TEXTRACT_RESULT_CACHE_PATH}/") if file_base_name in x]:    
            response = get_object_with_retry(BUCKET, f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt")
            text = response['Body'].read().decode()
            return text
        else:
            
            extractor = Textractor(region_name="us-west-2")
            # Asynchronous call, you will experience some wait time. Try caching results for better experience
            if "pdf" in ext:
                st.write("Asynchronous call, you may experience some wait time.")
                document = extractor.start_document_analysis(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],       
                save_image=False,   
                s3_output_path=f"s3://{BUCKET}/textract_output/"
            )
            #Synchronous call
            else:
                document = extractor.analyze_document(
                file_source=file,
                features=[TextractFeatures.LAYOUT,TextractFeatures.TABLES],  
                save_image=False,
            )
            config = TextLinearizationConfig(
            hide_figure_layout=False,   
            hide_header_layout=False,    
            table_prefix="<table>",
            table_suffix="</table>",
            )
            # Upload extracted content to s3
            S3.put_object(Body=document.get_text(config=config), Bucket=BUCKET, Key=f"{TEXTRACT_RESULT_CACHE_PATH}/{file_base_name}.txt") 
            return document.get_text(config=config)
    else:
        s3=boto3.resource("s3")
        match = re.match("s3://(.+?)/(.+)", file)
        if match:
            bucket_name = match.group(1)
            key = match.group(2)
        if "pdf" in ext:            
            pdf_bytes = io.BytesIO()            
            s3.Bucket(bucket_name).download_fileobj(key, pdf_bytes)
            # Read the PDF from the BytesIO object
            pdf_bytes.seek(0)                      
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(pdf_bytes)
            # Get the number of pages in the PDF
            num_pages = len(pdf_reader.pages)
            # Extract text from each page
            text = ''
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        else:
            img_bytes = io.BytesIO()
            s3.Bucket(bucket_name).download_fileobj(key, img_bytes)
            img_bytes.seek(0)
            image = Image.open(img_bytes)
            text = pytesseract.image_to_string(image)
        return text    

def strip_newline(cell):
    return str(cell).strip()

def table_parser_openpyxl(file):
    # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        # Load workbook
        wb = openpyxl.load_workbook(xlsx_buffer)    
        all_sheets_string=""
        # Iterate over each sheet in the workbook
        for sheet_name in wb.sheetnames:
            # all_sheets_name.append(sheet_name)
            worksheet = wb[sheet_name]

            all_merged_cell_ranges: list[CellRange] = list(
                worksheet.merged_cells.ranges
            )
            for merged_cell_range in all_merged_cell_ranges:
                merged_cell: Cell = merged_cell_range.start_cell
                worksheet.unmerge_cells(range_string=merged_cell_range.coord)
                for row_index, col_index in merged_cell_range.cells:
                    cell: Cell = worksheet.cell(row=row_index, column=col_index)
                    cell.value = merged_cell.value        
            # Convert sheet data to a DataFrame
            df = pd.DataFrame(worksheet.values)
            df = df.map(strip_newline)
            # Convert to string and tag by sheet name
            tabb=df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string+=f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def calamaine_excel_engine(file):
    # # Read from S3
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)
        obj = s3.get_object(Bucket=bucket_name, Key=key)    
        # Read Excel file from S3 into a buffer
        xlsx_buffer = io.BytesIO(obj['Body'].read())
        xlsx_buffer.seek(0)    
        all_sheets_string=""
        # Load the Excel file
        workbook = CalamineWorkbook.from_filelike(xlsx_buffer)
        # Iterate over each sheet in the workbook
        for sheet_name in workbook.sheet_names:
            # Get the sheet by name
            sheet = workbook.get_sheet_by_name(sheet_name)
            df = pd.DataFrame(sheet.to_python(skip_empty_area=False))
            df = df.map(strip_newline)
            tabb=df.to_csv(sep=CSV_SEPERATOR, index=False, header=0)
            all_sheets_string+=f'<{sheet_name}>\n{tabb}\n</{sheet_name}>\n'
        return all_sheets_string
    else:
        raise Exception(f"{file} not formatted as an S3 path")

def table_parser_utills(file):
    try:
        response= table_parser_openpyxl(file)
        if response:
            return response
        else:
            return calamaine_excel_engine(file)        
    except Exception as e:
        try:
            return calamaine_excel_engine(file)
        except Exception as e:
            raise Exception(str(e))


def put_db(params,messages):
    """Store long term chat history in DynamoDB"""    
    chat_item = {
        "UserId": st.session_state['userid'], # user id
        "SessionId": params["session_id"], # User session id
        "messages": [messages],  # 'messages' is a list of dictionaries
        "time":messages['time']
    }

    existing_item = DYNAMODB.Table(DYNAMODB_TABLE).get_item(Key={"UserId": st.session_state['userid'], "SessionId":params["session_id"]})
    if "Item" in existing_item:
        existing_messages = existing_item["Item"]["messages"]
        chat_item["messages"] = existing_messages + [messages]
    response = DYNAMODB.Table(DYNAMODB_TABLE).put_item(
        Item=chat_item
    )
    
    
def get_chat_history_db(params,cutoff,claude3):
    current_chat, chat_hist=[],[]
    if params['chat_histories']: 
        chat_hist=params['chat_histories'][-cutoff:]              
        for d in chat_hist:
            if d['image'] and claude3 and LOAD_DOC_IN_ALL_CHAT_CONVO:
                content=[]
                for img in d['image']:
                    s3 = boto3.client('s3')
                    match = re.match("s3://(.+?)/(.+)", img)
                    image_name=os.path.basename(img)
                    _,ext=os.path.splitext(image_name)
                    if "jpg" in ext: ext=".jpeg"                        
                    if match:
                        bucket_name = match.group(1)
                        key = match.group(2)    
                        obj = s3.get_object(Bucket=bucket_name, Key=key)
                        base_64_encoded_data = base64.b64encode(obj['Body'].read())
                        base64_string = base_64_encoded_data.decode('utf-8')                        
                    content.extend([{"type":"text","text":image_name},{
                      "type": "image",
                      "source": {
                        "type": "base64",
                        "media_type": f"image/{ext.lower().replace('.','')}",
                        "data": base64_string
                      }
                    }])
                content.extend([{"type":"text","text":d['user']}])
                current_chat.append({'role': 'user', 'content': content})
            elif d['document'] and LOAD_DOC_IN_ALL_CHAT_CONVO:
                doc='Here are the documents:\n'
                for docs in d['document']:
                    uploads=handle_doc_upload_or_s3(docs)
                    doc_name=os.path.basename(docs)
                    doc+=f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                if not claude3 and d["image"]:
                    for docs in d['image']:
                        uploads=handle_doc_upload_or_s3(docs)
                        doc_name=os.path.basename(docs)
                        doc+=f"<{doc_name}>\n{uploads}\n</{doc_name}>\n"
                current_chat.append({'role': 'user', 'content': [{"type":"text","text":doc+d['user']}]})
            else:
                current_chat.append({'role': 'user', 'content': [{"type":"text","text":d['user']}]})
            current_chat.append({'role': 'assistant', 'content': d['assistant']})  
    else:
        chat_hist=[]
    return current_chat, chat_hist

  
def get_s3_keys(prefix):
    """list all keys in an s3 path"""
    s3 = boto3.client('s3')
    keys = []
    next_token = None
    while True:
        if next_token:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix, ContinuationToken=next_token)
        else:
            response = s3.list_objects_v2(Bucket=BUCKET, Prefix=prefix)
        if "Contents" in response:
            for obj in response['Contents']:
                key = obj['Key']
                name = key[len(prefix):]
                keys.append(name)
        if "NextContinuationToken" in response:
            next_token = response["NextContinuationToken"]
        else:
            break
    return keys

def get_s3_obj_from_bucket_(file):
    s3 = boto3.client('s3')
    match = re.match("s3://(.+?)/(.+)", file)
    if match:
        bucket_name = match.group(1)
        key = match.group(2)    
        obj = s3.get_object(Bucket=bucket_name, Key=key)  
    return obj

def put_obj_in_s3_bucket_(docs):
    file_name=os.path.basename(docs.name)
    file_path=f"{S3_DOC_CACHE_PATH}/{file_name}"
    S3.put_object(Body=docs.read(),Bucket= BUCKET, Key=file_path)
    return f"s3://{BUCKET}/{file_path}"


def bedrock_streemer(params,response, handler):
    stream = response.get('body')
    answer = ""    
    if stream:
        for event in stream:
            chunk = event.get('chunk')
            if  chunk:
                chunk_obj = json.loads(chunk.get('bytes').decode())
                if "delta" in chunk_obj:                    
                    delta = chunk_obj['delta']
                    if "text" in delta:
                        text=delta['text'] 
                        # st.write(text, end="")                        
                        answer+=str(text)       
                        handler.markdown(answer.replace("$","USD ").replace("%", " percent"))
                        
                if "amazon-bedrock-invocationMetrics" in chunk_obj:
                    st.session_state['input_token'] = chunk_obj['amazon-bedrock-invocationMetrics']['inputTokenCount']
                    st.session_state['output_token'] =chunk_obj['amazon-bedrock-invocationMetrics']['outputTokenCount']
                    pricing=st.session_state['input_token']*pricing_file[f"anthropic.{params['model']}"]["input"]+st.session_state['output_token'] *pricing_file[f"anthropic.{params['model']}"]["output"]
                    st.session_state['cost']+=pricing             
    return answer

def bedrock_claude_(params,chat_history,system_message, prompt,model_id,image_path=None, handler=None):
    content=[]
    if image_path:       
        if not isinstance(image_path, list):
            image_path=[image_path]      
        for img in image_path:
            s3 = boto3.client('s3')
            match = re.match("s3://(.+?)/(.+)", img)
            image_name=os.path.basename(img)
            _,ext=os.path.splitext(image_name)
            if "jpg" in ext.lower(): ext=".jpeg"                        
            if match:
                bucket_name = match.group(1)
                key = match.group(2)    
                obj = s3.get_object(Bucket=bucket_name, Key=key)
                base_64_encoded_data = base64.b64encode(obj['Body'].read())
                base64_string = base_64_encoded_data.decode('utf-8')
            content.extend([{"type":"text","text":image_name},{
              "type": "image",
              "source": {
                "type": "base64",
                "media_type": f"image/{ext.lower().replace('.','')}",
                "data": base64_string
              }
            }])
    content.append({
        "type": "text",
        "text": prompt
            })
    chat_history.append({"role": "user",
            "content": content})
    # print(system_message)
    prompt = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 2500,
        "temperature": 0.5,
        "system":system_message,
        "messages": chat_history
    }
    
    prompt = json.dumps(prompt)
    response = bedrock_runtime.invoke_model_with_response_stream(body=prompt, modelId=model_id, accept="application/json", contentType="application/json")
    answer=bedrock_streemer(params,response, handler) 
    return answer

def _invoke_bedrock_with_retries(params,current_chat, chat_template, question, model_id, image_path, handler):
    max_retries = 10
    backoff_base = 2
    max_backoff = 3  # Maximum backoff time in seconds
    retries = 0

    while True:
        try:
            response = bedrock_claude_(params,current_chat, chat_template, question, model_id, image_path, handler)
            return response
        except ClientError as e:
            if e.response['Error']['Code'] == 'ThrottlingException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'ModelStreamErrorException':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            elif e.response['Error']['Code'] == 'EventStreamError':
                if retries < max_retries:
                    # Throttling, exponential backoff
                    sleep_time = min(max_backoff, backoff_base ** retries + random.uniform(0, 1))
                    time.sleep(sleep_time)
                    retries += 1
                else:
                    raise e
            else:
                # Some other API error, rethrow
                raise
                

def get_session_ids_by_user(table_name, user_id):
    """
    Get Session Ids and corresponding top message for a user to populate the chat history drop down on the front end
    """
    if DYNAMODB_TABLE:
        table = DYNAMODB.Table(table_name)
        message_list={}
        session_ids = []
        args = {
            'KeyConditionExpression': Key('UserId').eq(user_id)
        }
        while True:
            response = table.query(**args)
            session_ids.extend([item['SessionId'] for item in response['Items']])
            if 'LastEvaluatedKey' not in response:
                break
            args['ExclusiveStartKey'] = response['LastEvaluatedKey']

        for session_id in session_ids:
            try:
                message_list[session_id]=DYNAMODB.Table(table_name).get_item(Key={"UserId": user_id, "SessionId":session_id})['Item']['messages'][0]['user']
            except Exception as e:
                print(e)
                pass
    else:
        try:
            message_list={}
            # Read the existing JSON data from the file
            with open(LOCAL_CHAT_FILE_NAME, "r", encoding='utf-8') as file:
                existing_data = json.load(file)
            for session_id in existing_data:
                message_list[session_id]=existing_data[session_id][0]['user']
            
        except FileNotFoundError:
            # If the file doesn't exist, initialize an empty list
            message_list = {}
    return message_list


def query_llm(params, handler):
    """
    Function takes a user query and a uploaded document. Caches documents in S3
    passing a document is optional
    """  

    if not isinstance(params['upload_doc'], list):
        raise TypeError("documents must be in a list format")        
    # use CLAUDE 3 Model
    claude3=True
    model='anthropic.claude-3-sonnet-20240229-v1:0' 
    current_chat,chat_hist=get_chat_history_db(params, CHAT_HISTORY_LENGTH,claude3)

    ## prompt template for when a user uploads a doc
    doc_path=[]
    image_path=[]
    full_doc_path=[]
    doc=""
    if params['upload_doc']:  
        doc='I have provided documents and/or images.\n'
        for ids,docs in enumerate(params['upload_doc']):
            file_name=docs.name
            _,extensions=os.path.splitext(file_name)
            docs=put_obj_in_s3_bucket_(docs)
            full_doc_path.append(docs)
            if extensions.lower() in [".jpg",".jpeg",".png",".gif",".webp"] and claude3:
                image_path.append(docs)
                continue
        
        doc_path = [item for item in full_doc_path if item not in image_path]
        errors, result_string=process_files(doc_path)    
        if errors:
            st.error(errors)
        doc+= result_string
        with open("prompt/doc_chat.txt","r", encoding="utf-8") as f:
            chat_template=f.read()  
    else:        
        # Chat template for open ended query
        with open("prompt/chat.txt","r",encoding="utf-8") as f:
            chat_template=f.read()

    response=_invoke_bedrock_with_retries(params,current_chat, chat_template, doc+params['question'], model, image_path, handler)
    # log the following items to dynamodb
    chat_history={"user":params['question'],
    "assistant":response,
    "image":image_path,
    "document":doc_path,
    "modelID":model,
    "time":str(time.time()),
    "input_token":round(st.session_state['input_token']) ,
    "output_token":round(st.session_state['output_token'])} 
    #store convsation memory and user other items in DynamoDB table
    if DYNAMODB_TABLE:
        put_db(params,chat_history)
    # use local memory for storage
    else:
        save_chat_local(LOCAL_CHAT_FILE_NAME,[chat_history], params["session_id"])  
    return response


def get_chat_historie_for_streamlit(params):
    """
    This function retrieves chat history stored in a dynamoDB table partitioned by a userID and sorted by a SessionID
    """
    if DYNAMODB_TABLE:
        chat_histories = DYNAMODB.Table(DYNAMODB_TABLE).get_item(Key={"UserId": st.session_state['userid'], "SessionId":params["session_id"]})
        if "Item" in chat_histories:
            chat_histories=chat_histories['Item']['messages'] 
        else:
            chat_histories=[]
    else:
        chat_histories=load_chat_local(LOCAL_CHAT_FILE_NAME,params["session_id"])         

# Constructing the desired list of dictionaries
    formatted_data = []   
    if chat_histories:
        for entry in chat_histories:           
            image_files=[os.path.basename(x) for x in entry.get('image', [])]
            doc_files=[os.path.basename(x) for x in entry.get('document', [])]
            assistant_attachment = '\n\n'.join(image_files+doc_files)
            
            formatted_data.append({
                "role": "user",
                "content": entry["user"],
            })
            formatted_data.append({
                "role": "assistant",
                "content": entry["assistant"],
                "attachment": assistant_attachment
            })
    else:
        chat_histories=[]            
    return formatted_data,chat_histories



def get_key_from_value(dictionary, value):
    return next((key for key, val in dictionary.items() if val == value), None)

def chat_bedrock_(params):
    st.title('節能會議小助手💡')
    params['chat_histories']=[]
    if params["session_id"].strip():
        st.session_state.messages, params['chat_histories']=get_chat_historie_for_streamlit(params)
    for message in st.session_state.messages:

        with st.chat_message(message["role"]):   
            if "```" in message["content"]:
                st.markdown(message["content"],unsafe_allow_html=True )
            else:
                st.markdown(message["content"].replace("$", "\$"),unsafe_allow_html=True )
            if message["role"]=="assistant":
                if message["attachment"]:
                    with st.expander(label="**attachments**"):
                        st.markdown( message["attachment"])
        
    if prompt := st.chat_input("問我任何問題(๑•ૅω•´๑)"): 
        st.session_state.messages.append({"role": "user", "content": prompt})        
        with st.chat_message("user"):             
            st.markdown(prompt.replace("$", "\$"),unsafe_allow_html=True )
        with st.chat_message("assistant"): 
            message_placeholder = st.empty()
            time_now=time.time()            
            params["question"]=prompt
            answer=query_llm(params, message_placeholder)
            message_placeholder.markdown(answer.replace("$", "\$"),unsafe_allow_html=True )
            st.session_state.messages.append({"role": "assistant", "content": answer}) 
        st.rerun()
        
def app_sidebar():
    with st.sidebar:
        st.write("-----")
        model='claude-3-sonnet'
        params={"model":model} 
        user_sess_id=get_session_ids_by_user(DYNAMODB_TABLE, st.session_state['userid'])
        float_keys = {float(key): value for key, value in user_sess_id.items()}
        sorted_messages = sorted(float_keys.items(), reverse=True)      
        sorted_messages.insert(0, (float(st.session_state['user_sess']),"New Chat"))          
        st.session_state['chat_session_list'] = dict(sorted_messages)
        chat_items=st.selectbox("**對話階段：**",st.session_state['chat_session_list'].values(),key="chat_sessions")
        session_id=get_key_from_value(st.session_state['chat_session_list'], chat_items)   
        file = st.file_uploader('上傳檔案：', accept_multiple_files=True, help="pdf,csv,txt,png,jpg,xlsx,json,py doc format supported") 
        if file and LOAD_DOC_IN_ALL_CHAT_CONVO:
            st.warning('You have set **load-doc-in-chat-history** to true. For better performance, remove upload files (by clicking **X**) **AFTER** first query **RESPONSE** on uploaded files. See the README for more info', icon="⚠️")
        params={"model":model, "session_id":str(session_id), "chat_item":chat_items, "upload_doc":file }    
        st.session_state['count']=1
        return params
def main():
    params=app_sidebar()
    chat_bedrock_(params)
   
    
if __name__ == '__main__':
    main()   
    